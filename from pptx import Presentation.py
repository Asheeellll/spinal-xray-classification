from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    def add_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[5]  # Title Only Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        for line in content_text.split("\n"):
            p = text_frame.add_paragraph()
            p.text = line
            p.space_after = Pt(10)
            p.font.size = Pt(24)
            p.font.bold = False
            p.font.color.rgb = RGBColor(50, 50, 50)
        
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Character Scope and Importance in the Indian Knowledge System"
    subtitle.text = "A Brief Overview\nYour Name & Date"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Slides with custom styling
    add_slide("What is Character in the Indian Knowledge System?", "- Character (सदाचार) is a core value in ancient Indian teachings.\n- Rooted in Dharma (duty), Satya (truth), Ahimsa (non-violence), and Seva (service).\n- Emphasized in Vedas, Upanishads, Bhagavad Gita, and Buddhist & Jain philosophies.")
    
    add_slide("Scope of Character in Indian Knowledge System", "- Spiritual Development – Leads to inner peace.\n- Social Harmony – Encourages respect and justice.\n- Education & Gurukul System – Character-building was central.\n- Leadership & Governance – Followed by leaders like Rama, Ashoka, and Gandhi.")
    
    add_slide("Importance of Good Character in Indian Tradition", "- Moral Strength – Builds integrity and righteousness.\n- Karma & Rebirth – Actions shape future consequences.\n- Harmony with Nature – Respect for all living beings.\n- Role in Family & Society – Promotes respect and unity.")
    
    add_slide("How to Develop Good Character (Indian Perspective)", "- Practice Dharma – Follow righteous duties.\n- Truthfulness & Self-Discipline – Speak truth and control desires.\n- Bhakti & Seva (Devotion & Service) – Help others selflessly.\n- Meditation & Yoga – Develop inner strength and clarity.")
    
    add_slide("Conclusion", "- Character is the foundation of the Indian Knowledge System.\n- It influences personal growth, social values, and spiritual progress.\n- Living with Dharma, Truth, and Compassion leads to a fulfilling life.")
    
    # Save Presentation
    prs.save("Character_Scope_Indian_Knowledge_System.pptx")
    print("Presentation Created Successfully!")

# Run the function to create the PPT
create_presentation()
